"""
analyze_project_data.py  — v3
==============================
Reads an Excel file with project lessons-learned data and produces:
  1. An enriched Excel workbook  (Raw Data + insight sheets + embedded charts)
  2. A PowerPoint presentation   (title, agenda, chart slides, tables, takeaways)

NEW in v3:
  • N-gram analysis (bigrams & trigrams)
  • VADER sentiment analysis
  • Text length distribution histogram
  • TF-IDF distinctive keywords per category
  • Readability metrics (Flesch-Kincaid grade level)
  • Word co-occurrence network graph
  • LDA topic modeling (sklearn)
  • Text clustering  (TF-IDF + K-Means + PCA scatter)
  • Named entity recognition (NER) frequency (spaCy)
  • Sankey flow diagram (Type → Stage → Category, requires plotly + kaleido)
  • Semantic similarity heatmap (requires sentence-transformers)

Usage:
    python analyze_project_data.py  my_data.xlsx
    python analyze_project_data.py  my_data.xlsx  --sheet "Sheet2"  --out-dir ./results

Required Python packages:
    pip install pandas openpyxl python-pptx matplotlib

Optional (enables additional analyses):
    pip install nltk wordcloud pillow
    pip install scikit-learn
    pip install textstat
    pip install networkx
    pip install spacy && python -m spacy download en_core_web_sm
    pip install plotly kaleido
    pip install sentence-transformers

Column recognition is automatic — headings are matched case-insensitively.
"""

import sys, re, argparse, warnings
from collections import Counter
from pathlib import Path

import pandas as pd
import numpy as np
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.drawing.image import Image as XLImage

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

warnings.filterwarnings("ignore")

# ── Optional NLP (graceful fallback if not installed) ─────────────────────────
try:
    import nltk
    from nltk.corpus import stopwords as nltk_sw
    from nltk.stem import WordNetLemmatizer
    for _p in ["punkt", "stopwords", "wordnet", "omw-1.4", "punkt_tab", "vader_lexicon"]:
        try: nltk.download(_p, quiet=True)
        except: pass
    NLTK_OK = True
except ImportError:
    NLTK_OK = False

try:
    from wordcloud import WordCloud
    WC_OK = True
except ImportError:
    WC_OK = False

try:
    from sklearn.feature_extraction.text import CountVectorizer, TfidfVectorizer
    from sklearn.cluster import KMeans
    from sklearn.decomposition import PCA, LatentDirichletAllocation
    from sklearn.metrics.pairwise import cosine_similarity as sk_cosine
    SKLEARN_OK = True
except ImportError:
    SKLEARN_OK = False

try:
    import textstat
    TEXTSTAT_OK = True
except ImportError:
    TEXTSTAT_OK = False

try:
    import networkx as nx
    NX_OK = True
except ImportError:
    NX_OK = False

try:
    import spacy
    _nlp = spacy.load("en_core_web_sm")
    SPACY_OK = True
except (ImportError, OSError):
    SPACY_OK = False
    _nlp = None

try:
    import plotly.graph_objects as go
    import plotly.io as _pio
    PLOTLY_OK = True
except ImportError:
    PLOTLY_OK = False

try:
    from sentence_transformers import SentenceTransformer
    SBERT_OK = True
except ImportError:
    SBERT_OK = False

# ── Design constants ──────────────────────────────────────────────────────────
PAL = {
    "primary":     "1C4E80",
    "secondary":   "0091D5",
    "accent1":     "EA6A47",
    "accent2":     "A5D8DD",
    "accent3":     "488A99",
    "light":       "F0F4F8",
    "dark":        "1A2B3C",
    "threat":      "C0392B",
    "opportunity": "27AE60",
}
OPXL_HEX = {k: f"FF{v}" for k, v in PAL.items()}
PPTX_RGB = {k: RGBColor.from_string(v) for k, v in PAL.items()}

# ── Column aliases (fuzzy, case-insensitive) ──────────────────────────────────
ALIASES = {
    "project_name":      ["project name", "project"],
    "organization":      ["organization", "org", "department"],
    "project_type":      ["project type", "type"],
    "project_stage":     ["project stage", "stage", "phase"],
    "original_category": ["original category", "orig category", "orig cat"],
    "new_category":      ["new category", "revised category", "new cat"],
    "lesson_title":      ["lesson learned title", "lesson title", "title"],
    "threat_opp_type":   ["threat/opportunity", "threat or opportunity",
                          "risk type", "type of risk"],
    "description":       ["description of threat/opportunity", "description",
                          "threat description", "opportunity description", "desc"],
    "response":          ["response of threat/opportunity", "response",
                          "mitigation", "action"],
    "lessons_learned":   ["lessons learned", "lesson learned", "lessons",
                          "key takeaway", "takeaway"],
}
EXTRA_SW = {
    "project","projects","use","used","using","need","needs","ensure","must",
    "also","may","well","one","two","make","team","time","new","work","include",
    "process","provide","within","would","could","should","will","many","often",
    "however","therefore","thus","due","per","key","based","required",
    "high","low","good","important",
}

# =============================================================================
# DATA LOADING
# =============================================================================

def load_data(path: str, sheet: str = None) -> pd.DataFrame:
    xl = pd.ExcelFile(path)
    df = pd.read_excel(xl, sheet_name=sheet or xl.sheet_names[0], dtype=str).fillna("")
    df.columns = df.columns.str.strip()
    return df


def map_columns(df: pd.DataFrame) -> dict:
    lower = {c.lower().strip(): c for c in df.columns}
    result = {}
    for canon, aliases in ALIASES.items():
        for alias in aliases:
            if alias in lower:
                result[canon] = lower[alias]
                break
    return result


# =============================================================================
# TEXT UTILITIES
# =============================================================================

def build_stopwords() -> set:
    sw = EXTRA_SW.copy()
    if NLTK_OK:
        sw |= set(nltk_sw.words("english"))
    else:
        sw |= {
            "the","a","an","and","or","but","in","on","at","to","for","of",
            "with","by","from","is","are","was","were","be","been","being",
            "have","has","had","do","does","did","not","no","it","its",
            "they","them","their","we","our","you","your","he","she",
            "who","which","what","how","when","where","if","as",
            "this","that","these","those","then","here","there",
        }
    return sw


def tokenize(text: str, sw: set) -> list:
    tokens = re.sub(r"[^a-z\s]", " ", str(text).lower()).split()
    if NLTK_OK:
        lem = WordNetLemmatizer()
        tokens = [lem.lemmatize(t) for t in tokens]
    return [t for t in tokens if t not in sw and len(t) > 2]


def top_keywords(texts: list, sw: set, n: int = 20) -> list:
    c = Counter()
    for t in texts:
        c.update(tokenize(t, sw))
    return c.most_common(n)


def freq_table(df: pd.DataFrame, col: str, label: str) -> pd.DataFrame:
    counts = df[col].value_counts().reset_index()
    counts.columns = [label, "Count"]
    counts["Percentage"] = (counts["Count"] / counts["Count"].sum() * 100).round(1)
    return counts


def split_threats_opps(df: pd.DataFrame, col_map: dict):
    threats, opps = [], []
    dc = col_map.get("description", "")
    tc = col_map.get("threat_opp_type", "")
    if not dc:
        return threats, opps
    for _, row in df.iterrows():
        txt = str(row.get(dc, ""))
        lbl = str(row.get(tc, "")).lower() if tc else ""
        (opps if "opportunit" in lbl else threats).append(txt)
    return threats, opps


def stage_cat_pivot(df: pd.DataFrame, col_map: dict):
    sc = col_map.get("project_stage")
    cc = col_map.get("new_category") or col_map.get("original_category")
    if not sc or not cc:
        return None
    return pd.crosstab(df[sc], df[cc])


# =============================================================================
# ANALYTICS (NLP / ML — all optional)
# =============================================================================

def run_analytics(df: pd.DataFrame, col_map: dict, sw: set) -> dict:
    """
    Run all optional NLP/ML analyses.
    Returns a dict of computed results used by generate_charts() and build_excel().
    All analyses are optional — missing packages produce None values gracefully.
    """
    analytics = {}

    # Gather common text collections
    lc  = col_map.get("lessons_learned")
    dc  = col_map.get("description")
    rc  = col_map.get("response")
    tc  = col_map.get("lesson_title")
    cat_col  = col_map.get("new_category") or col_map.get("original_category")
    stage_col = col_map.get("project_stage")

    lesson_texts = df[lc].tolist() if lc else []
    all_text_cols = [c for c in [dc, rc, lc, tc] if c]
    all_texts = []
    for c in all_text_cols:
        all_texts += df[c].tolist()

    # 1. VADER Sentiment ───────────────────────────────────────────────────────
    if NLTK_OK and lc:
        try:
            from nltk.sentiment.vader import SentimentIntensityAnalyzer
            sia = SentimentIntensityAnalyzer()
            scores = df[lc].apply(lambda t: sia.polarity_scores(str(t)))
            sent_df = pd.DataFrame(scores.tolist())
            sent_df["label"] = sent_df["compound"].apply(
                lambda s: "Positive" if s >= 0.05 else ("Negative" if s <= -0.05 else "Neutral"))
            analytics["sentiment"] = sent_df
        except Exception as e:
            print(f"  [SKIP] VADER sentiment: {e}")
            analytics["sentiment"] = None
    else:
        analytics["sentiment"] = None

    # 2. N-gram Analysis ───────────────────────────────────────────────────────
    if SKLEARN_OK:
        analytics["ngrams"] = {}
        for field, col in [("description", dc), ("response", rc), ("lessons_learned", lc)]:
            if not col:
                continue
            texts = df[col].tolist()
            field_ngrams = {}
            for n, name in [(2, "bigrams"), (3, "trigrams")]:
                try:
                    vec = CountVectorizer(
                        ngram_range=(n, n), stop_words=list(sw),
                        max_features=500, min_df=1)
                    X = vec.fit_transform(texts)
                    counts = X.sum(axis=0).A1
                    names = vec.get_feature_names_out()
                    top = sorted(zip(names, counts.tolist()), key=lambda x: -x[1])[:15]
                    field_ngrams[name] = [(phrase, int(cnt)) for phrase, cnt in top]
                except Exception:
                    field_ngrams[name] = []
            analytics["ngrams"][field] = field_ngrams
    else:
        analytics["ngrams"] = {}

    # 3. TF-IDF Keywords per Category ─────────────────────────────────────────
    if SKLEARN_OK and cat_col and (lc or dc):
        text_col = lc or dc
        tfidf_cats = {}
        for cat in df[cat_col].unique():
            texts = df[df[cat_col] == cat][text_col].tolist()
            if len(texts) < 2:
                continue
            try:
                vec = TfidfVectorizer(stop_words=list(sw), max_features=300)
                X = vec.fit_transform(texts)
                scores = X.mean(axis=0).A1
                names = vec.get_feature_names_out()
                top = sorted(zip(names, scores.tolist()), key=lambda x: -x[1])[:10]
                tfidf_cats[str(cat)] = [(w, round(s, 4)) for w, s in top]
            except Exception:
                continue
        analytics["tfidf_categories"] = tfidf_cats
    else:
        analytics["tfidf_categories"] = {}

    # 4. Readability Metrics ───────────────────────────────────────────────────
    if TEXTSTAT_OK and lc:
        def _metrics(text):
            t = str(text)
            words = t.split()
            if len(words) < 5:
                return {"flesch_kincaid": None, "avg_word_length": None, "word_count": len(words)}
            try:
                return {
                    "flesch_kincaid":  textstat.flesch_kincaid_grade(t),
                    "avg_word_length": round(sum(len(w) for w in words) / len(words), 2),
                    "word_count":      len(words),
                }
            except Exception:
                return {"flesch_kincaid": None, "avg_word_length": None, "word_count": len(words)}
        analytics["readability"] = pd.DataFrame(df[lc].apply(_metrics).tolist())
    else:
        analytics["readability"] = None

    # 5. LDA Topic Modeling ────────────────────────────────────────────────────
    if SKLEARN_OK and len(lesson_texts) >= 10:
        try:
            n_topics = min(5, max(2, len(lesson_texts) // 10))
            vec = CountVectorizer(stop_words=list(sw), max_features=300, min_df=2)
            X = vec.fit_transform(lesson_texts)
            if X.shape[1] > 0:
                lda = LatentDirichletAllocation(
                    n_components=n_topics, random_state=42, max_iter=30)
                lda.fit(X)
                feat_names = vec.get_feature_names_out()
                topics = []
                for comp in lda.components_:
                    top_words = [feat_names[i] for i in comp.argsort()[:-11:-1]]
                    topics.append(top_words)
                analytics["topics"] = topics
            else:
                analytics["topics"] = None
        except Exception as e:
            print(f"  [SKIP] LDA topic modeling: {e}")
            analytics["topics"] = None
    else:
        analytics["topics"] = None

    # 6. Named Entity Recognition ─────────────────────────────────────────────
    if SPACY_OK and all_texts:
        try:
            ent_counter = Counter()
            for text in all_texts:
                doc = _nlp(str(text)[:500_000])
                for ent in doc.ents:
                    if ent.label_ in {"ORG", "PERSON", "GPE", "PRODUCT", "EVENT"}:
                        ent_counter[(ent.text.strip(), ent.label_)] += 1
            analytics["ner"] = ent_counter
        except Exception as e:
            print(f"  [SKIP] NER: {e}")
            analytics["ner"] = None
    else:
        analytics["ner"] = None

    # 7. Text Lengths ──────────────────────────────────────────────────────────
    analytics["lengths"] = df[lc].apply(lambda t: len(str(t).split())) if lc else None

    # 8. Groupby column for readability chart ─────────────────────────────────
    analytics["_cat_col"]   = cat_col
    analytics["_stage_col"] = stage_col

    return analytics


# =============================================================================
# CHART GENERATION
# =============================================================================

CHART_DIR = Path("/tmp/pld_charts")
CHART_DIR.mkdir(exist_ok=True)

MULTI_COLORS = [
    f"#{PAL['primary']}", f"#{PAL['secondary']}", f"#{PAL['accent1']}",
    f"#{PAL['accent2']}", f"#{PAL['accent3']}", "#7EC8E3", "#FFD166",
    "#EF476F", "#06D6A0", "#118AB2"
] * 4


def _setup(fig, ax):
    fig.patch.set_facecolor("#FAFAFA")
    ax.set_facecolor("#FAFAFA")
    ax.spines[["top", "right"]].set_visible(False)


def make_bar(labels, values, title, fname, color=None, horiz=False, figsize=(8, 4)) -> str:
    color = color or f"#{PAL['secondary']}"
    fig, ax = plt.subplots(figsize=figsize)
    _setup(fig, ax)
    pos = range(len(labels))
    if horiz:
        ax.barh(pos, values, color=color, edgecolor="white", linewidth=0.5)
        ax.set_yticks(pos)
        ax.set_yticklabels(labels, fontsize=9)
        ax.invert_yaxis()
        ax.set_xlabel("Count", fontsize=9)
        for i, v in enumerate(values):
            ax.text(v + max(values) * 0.01, i, str(v), va="center", fontsize=8)
    else:
        ax.bar(pos, values, color=color, edgecolor="white", linewidth=0.5)
        ax.set_xticks(pos)
        ax.set_xticklabels(labels, rotation=35, ha="right", fontsize=9)
        ax.set_ylabel("Count", fontsize=9)
        for i, v in enumerate(values):
            ax.text(i, v + max(values) * 0.01, str(v), ha="center", fontsize=8)
    ax.set_title(title, fontsize=11, fontweight="bold", color=f"#{PAL['dark']}", pad=10)
    plt.tight_layout()
    p = CHART_DIR / fname
    plt.savefig(p, dpi=150, bbox_inches="tight")
    plt.close()
    return str(p)


def make_pie(labels, values, title, fname) -> str:
    fig, ax = plt.subplots(figsize=(7, 5))
    fig.patch.set_facecolor("#FAFAFA")
    wedges, _, autotexts = ax.pie(
        values, colors=MULTI_COLORS[:len(values)], autopct="%1.1f%%",
        startangle=140, wedgeprops=dict(edgecolor="white", linewidth=1.5))
    for at in autotexts:
        at.set_fontsize(8)
    ax.legend(wedges, labels, loc="center left", bbox_to_anchor=(1, 0, 0.5, 1), fontsize=8)
    ax.set_title(title, fontsize=11, fontweight="bold", color=f"#{PAL['dark']}", pad=10)
    plt.tight_layout()
    p = CHART_DIR / fname
    plt.savefig(p, dpi=150, bbox_inches="tight")
    plt.close()
    return str(p)


def make_kw_chart(kw_list, title, fname, color=None) -> str:
    if not kw_list:
        return None
    labels = [k for k, _ in kw_list[:15]]
    values = [v for _, v in kw_list[:15]]
    return make_bar(labels, values, title, fname, color=color, horiz=True, figsize=(8, 5))


def make_wordcloud(texts, fname) -> str:
    if not WC_OK or not texts:
        return None
    sw = build_stopwords()
    combined = " ".join(" ".join(tokenize(t, sw)) for t in texts)
    if not combined.strip():
        return None
    wc = WordCloud(width=800, height=400, background_color="white",
                   colormap="Blues", max_words=80,
                   stopwords=sw, collocations=False).generate(combined)
    fig, ax = plt.subplots(figsize=(8, 4))
    ax.imshow(wc, interpolation="bilinear")
    ax.axis("off")
    plt.tight_layout(pad=0)
    p = CHART_DIR / fname
    plt.savefig(p, dpi=150, bbox_inches="tight")
    plt.close()
    return str(p)


def make_stacked_bar(pivot_df, title, fname) -> str:
    fig, ax = plt.subplots(figsize=(10, 5))
    _setup(fig, ax)
    bottom = np.zeros(len(pivot_df))
    for i, col in enumerate(pivot_df.columns):
        ax.bar(pivot_df.index, pivot_df[col], bottom=bottom,
               label=col, color=MULTI_COLORS[i % len(MULTI_COLORS)],
               edgecolor="white", linewidth=0.5)
        bottom += pivot_df[col].values
    ax.set_title(title, fontsize=11, fontweight="bold", color=f"#{PAL['dark']}", pad=10)
    ax.set_ylabel("Count", fontsize=9)
    ax.set_xticklabels(pivot_df.index, rotation=30, ha="right", fontsize=9)
    ax.legend(bbox_to_anchor=(1.01, 1), loc="upper left", fontsize=8)
    plt.tight_layout()
    p = CHART_DIR / fname
    plt.savefig(p, dpi=150, bbox_inches="tight")
    plt.close()
    return str(p)


# ── New chart functions (v3) ──────────────────────────────────────────────────

def make_ngram_chart(ngram_list, title, fname, color=None) -> str:
    """Horizontal bar for top bigrams or trigrams."""
    if not ngram_list:
        return None
    labels = [p for p, _ in ngram_list[:12]]
    values = [c for _, c in ngram_list[:12]]
    if not values or max(values) == 0:
        return None
    return make_bar(labels, values, title, fname,
                    color=color or f"#{PAL['accent2']}", horiz=True, figsize=(9, 5))


def make_sentiment_charts(sentiment_df, fname_pie, fname_hist) -> dict:
    """Sentiment distribution pie + VADER compound score histogram."""
    if sentiment_df is None or sentiment_df.empty:
        return {}
    charts = {}
    label_counts = sentiment_df["label"].value_counts()
    color_map = {
        "Positive": f"#{PAL['opportunity']}",
        "Neutral":  f"#{PAL['secondary']}",
        "Negative": f"#{PAL['threat']}",
    }
    colors = [color_map.get(l, f"#{PAL['accent2']}") for l in label_counts.index]

    # Pie
    fig, ax = plt.subplots(figsize=(7, 5))
    fig.patch.set_facecolor("#FAFAFA")
    wedges, _, autotexts = ax.pie(
        label_counts.values, colors=colors, autopct="%1.1f%%",
        startangle=140, wedgeprops=dict(edgecolor="white", linewidth=1.5))
    for at in autotexts:
        at.set_fontsize(9)
    ax.legend(wedges, label_counts.index, loc="center left",
              bbox_to_anchor=(1, 0, 0.5, 1), fontsize=9)
    ax.set_title("Sentiment Distribution (VADER)", fontsize=11,
                 fontweight="bold", color=f"#{PAL['dark']}", pad=10)
    plt.tight_layout()
    p = CHART_DIR / fname_pie
    plt.savefig(p, dpi=150, bbox_inches="tight")
    plt.close()
    charts["sentiment_pie"] = str(p)

    # Histogram
    fig, ax = plt.subplots(figsize=(8, 4))
    _setup(fig, ax)
    ax.hist(sentiment_df["compound"], bins=25,
            color=f"#{PAL['secondary']}", edgecolor="white", linewidth=0.5)
    ax.axvline( 0.05, color=f"#{PAL['opportunity']}", linestyle="--",
               linewidth=1.5, label="Positive threshold (+0.05)")
    ax.axvline(-0.05, color=f"#{PAL['threat']}", linestyle="--",
               linewidth=1.5, label="Negative threshold (−0.05)")
    ax.set_xlabel("Compound Score", fontsize=9)
    ax.set_ylabel("Count", fontsize=9)
    ax.set_title("Sentiment Score Distribution (VADER Compound)",
                 fontsize=11, fontweight="bold", color=f"#{PAL['dark']}", pad=10)
    ax.legend(fontsize=9)
    plt.tight_layout()
    p = CHART_DIR / fname_hist
    plt.savefig(p, dpi=150, bbox_inches="tight")
    plt.close()
    charts["sentiment_hist"] = str(p)
    return charts


def make_length_hist(lengths, title, fname) -> str:
    """Histogram of word counts per record."""
    if lengths is None or len(lengths) == 0:
        return None
    fig, ax = plt.subplots(figsize=(8, 4))
    _setup(fig, ax)
    ax.hist(lengths, bins=min(30, max(len(lengths) // 3, 5)),
            color=f"#{PAL['accent2']}", edgecolor="white", linewidth=0.5)
    mean_val = lengths.mean()
    ax.axvline(mean_val, color=f"#{PAL['accent1']}", linestyle="--",
               linewidth=1.5, label=f"Mean: {mean_val:.1f} words")
    ax.set_xlabel("Word Count", fontsize=9)
    ax.set_ylabel("Number of Records", fontsize=9)
    ax.set_title(title, fontsize=11, fontweight="bold", color=f"#{PAL['dark']}", pad=10)
    ax.legend(fontsize=9)
    plt.tight_layout()
    p = CHART_DIR / fname
    plt.savefig(p, dpi=150, bbox_inches="tight")
    plt.close()
    return str(p)


def make_readability_chart(df, col_map, read_df, fname) -> str:
    """Avg Flesch-Kincaid grade level grouped by stage or category."""
    if read_df is None:
        return None
    group_col = col_map.get("project_stage") or col_map.get("new_category") or col_map.get("original_category")
    if not group_col:
        return None
    try:
        combined = pd.concat([
            df[[group_col]].reset_index(drop=True),
            read_df[["flesch_kincaid"]].reset_index(drop=True),
        ], axis=1)
        avg = combined.groupby(group_col)["flesch_kincaid"].mean().dropna().sort_values()
        if len(avg) < 2:
            return None
        label = group_col.replace("_", " ").title()
        return make_bar(
            avg.index.tolist(), avg.values.round(1).tolist(),
            f"Avg. Readability (Flesch-Kincaid Grade) by {label}",
            fname, color=f"#{PAL['accent1']}", horiz=True, figsize=(9, 5))
    except Exception as e:
        print(f"  [SKIP] Readability chart: {e}")
        return None


def make_cooccurrence_network(texts, sw, fname, top_n=40) -> str:
    """Word co-occurrence network saved as PNG (requires networkx)."""
    if not NX_OK or not texts:
        return None
    try:
        pair_counts = Counter()
        for text in texts:
            tokens = list(dict.fromkeys(tokenize(text, sw)))[:15]
            for i in range(len(tokens)):
                for j in range(i + 1, len(tokens)):
                    pair_counts[tuple(sorted([tokens[i], tokens[j]]))] += 1

        top_pairs = pair_counts.most_common(top_n)
        if len(top_pairs) < 5:
            return None

        G = nx.Graph()
        for (w1, w2), cnt in top_pairs:
            G.add_edge(w1, w2, weight=cnt)

        fig, ax = plt.subplots(figsize=(10, 8))
        fig.patch.set_facecolor("#FAFAFA")
        ax.set_facecolor("#FAFAFA")

        pos = nx.spring_layout(G, seed=42, k=2.0)
        weights = [G[u][v]["weight"] for u, v in G.edges()]
        max_w = max(weights) if weights else 1
        node_sizes = [300 + G.degree(n) * 80 for n in G.nodes()]

        nx.draw_networkx_edges(G, pos, width=[w / max_w * 3 for w in weights],
                               alpha=0.35, edge_color=f"#{PAL['accent2']}", ax=ax)
        nx.draw_networkx_nodes(G, pos, node_size=node_sizes,
                               node_color=f"#{PAL['secondary']}", alpha=0.9, ax=ax)
        nx.draw_networkx_labels(G, pos, font_size=7, font_color="white",
                                font_weight="bold", ax=ax)
        ax.set_title("Word Co-occurrence Network", fontsize=11,
                     fontweight="bold", color=f"#{PAL['dark']}", pad=10)
        ax.axis("off")
        plt.tight_layout()
        p = CHART_DIR / fname
        plt.savefig(p, dpi=150, bbox_inches="tight")
        plt.close()
        return str(p)
    except Exception as e:
        print(f"  [SKIP] Co-occurrence network: {e}")
        return None


def make_topic_chart(topics, fname) -> str:
    """Side-by-side panels showing top terms for each LDA topic."""
    if not topics:
        return None
    try:
        n = len(topics)
        fig, axes = plt.subplots(1, n, figsize=(max(3.5 * n, 10), 5))
        if n == 1:
            axes = [axes]
        fig.patch.set_facecolor("#FAFAFA")
        for i, (ax, words) in enumerate(zip(axes, topics)):
            ax.set_facecolor("#FAFAFA")
            scores = list(range(len(words), 0, -1))
            ax.barh(range(len(words)), scores,
                    color=MULTI_COLORS[i % len(MULTI_COLORS)], edgecolor="white")
            ax.set_yticks(range(len(words)))
            ax.set_yticklabels(words, fontsize=8)
            ax.set_title(f"Topic {i + 1}", fontweight="bold", fontsize=10,
                         color=f"#{PAL['dark']}")
            ax.set_xticks([])
            ax.spines[["top", "right", "bottom", "left"]].set_visible(False)
        fig.suptitle("LDA Topic Modeling — Top Terms per Topic",
                     fontsize=12, fontweight="bold", color=f"#{PAL['dark']}")
        plt.tight_layout()
        p = CHART_DIR / fname
        plt.savefig(p, dpi=150, bbox_inches="tight")
        plt.close()
        return str(p)
    except Exception as e:
        print(f"  [SKIP] Topic chart: {e}")
        return None


def make_cluster_scatter(texts, fname, n_clusters=4) -> str:
    """TF-IDF + K-Means + PCA 2D scatter of document clusters (requires sklearn)."""
    if not SKLEARN_OK or not texts or len(texts) < n_clusters * 3:
        return None
    try:
        clean = [re.sub(r"[^a-z\s]", " ", str(t).lower()) for t in texts]
        vec = TfidfVectorizer(max_features=500, min_df=1)
        X = vec.fit_transform(clean)
        if X.shape[0] < n_clusters:
            return None
        kmeans = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
        cluster_labels = kmeans.fit_predict(X)
        pca = PCA(n_components=2, random_state=42)
        coords = pca.fit_transform(X.toarray())

        fig, ax = plt.subplots(figsize=(8, 6))
        _setup(fig, ax)
        for ci in range(n_clusters):
            mask = cluster_labels == ci
            ax.scatter(coords[mask, 0], coords[mask, 1],
                       c=MULTI_COLORS[ci], label=f"Cluster {ci + 1}",
                       alpha=0.7, s=60, edgecolors="white", linewidth=0.5)
        ax.set_xlabel(
            f"PC1 ({pca.explained_variance_ratio_[0] * 100:.1f}% variance)", fontsize=9)
        ax.set_ylabel(
            f"PC2 ({pca.explained_variance_ratio_[1] * 100:.1f}% variance)", fontsize=9)
        ax.set_title("Document Clusters (TF-IDF + K-Means + PCA)",
                     fontsize=11, fontweight="bold", color=f"#{PAL['dark']}", pad=10)
        ax.legend(fontsize=9)
        plt.tight_layout()
        p = CHART_DIR / fname
        plt.savefig(p, dpi=150, bbox_inches="tight")
        plt.close()
        return str(p)
    except Exception as e:
        print(f"  [SKIP] Cluster scatter: {e}")
        return None


def make_ner_chart(ent_counter, fname) -> str:
    """Horizontal bar chart of top named entities (requires spaCy)."""
    if not ent_counter:
        return None
    top = ent_counter.most_common(15)
    if not top:
        return None
    labels = [f"{ent} ({lbl})" for (ent, lbl), _ in top]
    values = [cnt for _, cnt in top]
    return make_bar(labels, values, "Top Named Entities (NER)",
                    fname, color=f"#{PAL['accent3']}", horiz=True, figsize=(9, 5))


def make_sankey(df, col_map, fname) -> str:
    """Sankey flow diagram: Type → Stage → Category (requires plotly + kaleido)."""
    if not PLOTLY_OK:
        return None
    cols = [col_map.get(f) for f in ["project_type", "project_stage", "new_category"]
            if col_map.get(f)]
    if len(cols) < 2:
        return None
    try:
        all_labels, label_idx = [], {}
        for col in cols:
            for val in df[col].dropna().unique():
                key = f"{col}::{val}"
                if key not in label_idx:
                    label_idx[key] = len(all_labels)
                    all_labels.append(str(val))

        sources, targets, values = [], [], []
        for i in range(len(cols) - 1):
            grp = df.groupby([cols[i], cols[i + 1]]).size().reset_index(name="n")
            for _, row in grp.iterrows():
                s = label_idx.get(f"{cols[i]}::{row[cols[i]]}")
                t = label_idx.get(f"{cols[i + 1]}::{row[cols[i + 1]]}")
                if s is not None and t is not None:
                    sources.append(s)
                    targets.append(t)
                    values.append(int(row["n"]))

        if not sources:
            return None

        field_names = [c.replace("_", " ").title() for c in
                       ["Project Type", "Stage", "Category"][:len(cols)]]
        title = "Flow: " + " → ".join(field_names)

        fig = go.Figure(go.Sankey(
            node=dict(label=all_labels, pad=15, thickness=20,
                      color=f"#{PAL['secondary']}"),
            link=dict(source=sources, target=targets, value=values,
                      color="rgba(0,145,213,0.25)"),
        ))
        fig.update_layout(
            title_text=title, font_size=11,
            paper_bgcolor="#FAFAFA", width=1200, height=600)
        p = CHART_DIR / fname
        fig.write_image(str(p), scale=2)
        return str(p)
    except Exception as e:
        print(f"  [SKIP] Sankey diagram: {e}")
        return None


def make_similarity_heatmap(texts, row_labels, fname, max_docs=60) -> str:
    """Cosine similarity heatmap via sentence-transformers embeddings."""
    if not SBERT_OK or len(texts) < 3:
        return None
    try:
        texts      = [str(t) for t in texts[:max_docs]]
        row_labels = [str(l)[:28] for l in row_labels[:max_docs]]
        model      = SentenceTransformer("all-MiniLM-L6-v2")
        embeddings = model.encode(texts, show_progress_bar=False)
        sim_matrix = sk_cosine(embeddings)

        size = max(8, len(texts) * 0.25)
        fig, ax = plt.subplots(figsize=(size, size * 0.85))
        im = ax.imshow(sim_matrix, cmap="Blues", aspect="auto", vmin=0, vmax=1)
        plt.colorbar(im, ax=ax, shrink=0.8, label="Cosine Similarity")
        fs = max(5, 9 - len(texts) // 10)
        ax.set_xticks(range(len(row_labels)))
        ax.set_yticks(range(len(row_labels)))
        ax.set_xticklabels(row_labels, rotation=90, fontsize=fs)
        ax.set_yticklabels(row_labels, fontsize=fs)
        ax.set_title("Semantic Similarity Matrix (Sentence Embeddings)",
                     fontsize=11, fontweight="bold", color=f"#{PAL['dark']}", pad=10)
        plt.tight_layout()
        p = CHART_DIR / fname
        plt.savefig(p, dpi=150, bbox_inches="tight")
        plt.close()
        return str(p)
    except Exception as e:
        print(f"  [SKIP] Similarity heatmap: {e}")
        return None


def generate_charts(df: pd.DataFrame, col_map: dict, sw: set, analytics: dict) -> dict:
    charts = {}

    # ── Original charts ──────────────────────────────────────────────────────
    for field, fname, color in [
        ("project_type",  "type_bar.png",     f"#{PAL['primary']}"),
        ("project_stage", "stage_bar.png",    f"#{PAL['secondary']}"),
        ("new_category",  "category_bar.png", f"#{PAL['accent1']}"),
        ("organization",  "org_bar.png",      f"#{PAL['accent3']}"),
    ]:
        col = col_map.get(field)
        if not col:
            continue
        df_f   = freq_table(df, col, "")
        labels = df_f.iloc[:, 0].tolist()[:12]
        values = df_f["Count"].tolist()[:12]
        title  = f"Distribution by {field.replace('_', ' ').title()}"
        horiz  = max((len(str(l)) for l in labels), default=0) > 12
        charts[fname.replace(".png", "")] = make_bar(
            labels, values, title, fname, color=color, horiz=horiz)

    to_col = col_map.get("threat_opp_type")
    if to_col:
        df_to = freq_table(df, to_col, "")
        if len(df_to) >= 1:
            charts["threat_opp_pie"] = make_pie(
                df_to.iloc[:, 0].tolist(), df_to["Count"].tolist(),
                "Threat vs. Opportunity Split", "threat_opp_pie.png")

    all_text_cols = [col_map.get(f) for f in
                     ["description", "response", "lessons_learned", "lesson_title"]
                     if col_map.get(f)]
    all_texts = []
    for c in all_text_cols:
        all_texts += df[c].tolist()

    if all_texts:
        kws = top_keywords(all_texts, sw, 15)
        charts["keyword_all"] = make_kw_chart(
            kws, "Top Keywords — All Text", "keyword_all.png", f"#{PAL['primary']}")

    lc = col_map.get("lessons_learned")
    if lc:
        kws = top_keywords(df[lc].tolist(), sw, 15)
        charts["keyword_lessons"] = make_kw_chart(
            kws, "Top Keywords — Lessons Learned",
            "keyword_lessons.png", f"#{PAL['accent3']}")

    t_texts, o_texts = split_threats_opps(df, col_map)
    if t_texts:
        charts["threat_kw"] = make_kw_chart(
            top_keywords(t_texts, sw, 15),
            "Top Keywords — Threats", "threat_kw.png", f"#{PAL['threat']}")
    if o_texts:
        charts["opp_kw"] = make_kw_chart(
            top_keywords(o_texts, sw, 15),
            "Top Keywords — Opportunities", "opp_kw.png", f"#{PAL['opportunity']}")

    if all_texts:
        charts["wordcloud"] = make_wordcloud(all_texts, "wordcloud.png")

    pivot = stage_cat_pivot(df, col_map)
    if pivot is not None and not pivot.empty:
        charts["stage_cat_stack"] = make_stacked_bar(
            pivot, "Lessons by Stage & Category", "stage_cat_stack.png")

    # ── New v3 charts ────────────────────────────────────────────────────────

    # N-gram charts (bigrams & trigrams per text field)
    ngrams = analytics.get("ngrams", {})
    for field, field_ngrams in ngrams.items():
        for ng_type, ng_data in field_ngrams.items():
            if ng_data:
                label = field.replace("_", " ").title()
                title = f"Top {ng_type.title()} — {label}"
                fname = f"{field}_{ng_type}.png"
                charts[f"{field}_{ng_type}"] = make_ngram_chart(ng_data, title, fname)

    # Sentiment charts
    sent_df = analytics.get("sentiment")
    if sent_df is not None:
        charts.update(make_sentiment_charts(
            sent_df, "sentiment_pie.png", "sentiment_hist.png"))

    # Text length histogram
    lengths = analytics.get("lengths")
    if lengths is not None:
        charts["length_hist"] = make_length_hist(
            lengths, "Lessons Learned — Word Count Distribution", "length_hist.png")

    # Readability chart
    read_df = analytics.get("readability")
    if read_df is not None:
        charts["readability"] = make_readability_chart(
            df, col_map, read_df, "readability.png")

    # Word co-occurrence network
    if all_texts:
        charts["cooccurrence"] = make_cooccurrence_network(
            all_texts, sw, "cooccurrence.png")

    # LDA topic model chart
    topics = analytics.get("topics")
    if topics:
        charts["topic_model"] = make_topic_chart(topics, "topic_model.png")

    # Text cluster scatter plot
    if lc:
        n_clusters = min(5, max(2, len(df) // 15))
        charts["cluster_scatter"] = make_cluster_scatter(
            df[lc].tolist(), "cluster_scatter.png", n_clusters=n_clusters)

    # NER chart
    ner = analytics.get("ner")
    if ner:
        charts["ner_chart"] = make_ner_chart(ner, "ner_chart.png")

    # Sankey flow diagram
    charts["sankey"] = make_sankey(df, col_map, "sankey.png")

    # Semantic similarity heatmap
    tc = col_map.get("lesson_title")
    if lc:
        row_labels = df[tc].tolist() if tc else df[lc].apply(
            lambda t: str(t)[:30]).tolist()
        charts["similarity"] = make_similarity_heatmap(
            df[lc].tolist(), row_labels, "similarity.png")

    print(f"  {sum(1 for v in charts.values() if v)} chart(s) generated.")
    return charts


# =============================================================================
# EXCEL BUILDER
# =============================================================================

def _fill(c):
    return PatternFill("solid", fgColor=c)

def _font(bold=False, color="FF000000", size=10):
    return Font(bold=bold, color=color, size=size, name="Calibri")

def _align(h="left", v="center", wrap=False):
    return Alignment(horizontal=h, vertical=v, wrap_text=wrap)

def _bdr():
    s = Side(style="thin", color="FFD0D0D0")
    return Border(left=s, right=s, top=s, bottom=s)


def style_header(ws, row, ncols, bg=None, fg="FFFFFFFF"):
    bg = bg or OPXL_HEX["primary"]
    for c in range(1, ncols + 1):
        cell = ws.cell(row=row, column=c)
        cell.fill = _fill(bg)
        cell.font = _font(True, fg, 10)
        cell.alignment = _align("center")
        cell.border = _bdr()


def style_data(ws, r0, r1, ncols):
    for r in range(r0, r1 + 1):
        bg = OPXL_HEX["light"] if r % 2 == 0 else "FFFFFFFF"
        for c in range(1, ncols + 1):
            cell = ws.cell(row=r, column=c)
            cell.fill = _fill(bg)
            cell.font = _font(size=10)
            cell.alignment = _align(wrap=True)
            cell.border = _bdr()


def write_freq_block(ws, df_f, start_row, title="") -> int:
    if title:
        ws.cell(row=start_row, column=1).value = title
        ws.cell(row=start_row, column=1).font = _font(True, OPXL_HEX["primary"], 12)
        start_row += 1
    for ci, col in enumerate(df_f.columns, 1):
        ws.cell(row=start_row, column=ci).value = col
    style_header(ws, start_row, len(df_f.columns))
    for ri, row in df_f.iterrows():
        for ci, val in enumerate(row, 1):
            ws.cell(row=start_row + ri + 1, column=ci).value = val
    style_data(ws, start_row + 1, start_row + len(df_f), len(df_f.columns))
    return start_row + len(df_f) + 2


def embed_image(ws, path, anchor, w_cm=14, h_cm=7):
    if not path or not Path(path).exists():
        return
    img = XLImage(path)
    img.width  = int(w_cm * 37.795)
    img.height = int(h_cm * 37.795)
    ws.add_image(img, anchor)


def build_excel(df: pd.DataFrame, col_map: dict, charts: dict,
                sw: set, analytics: dict, out_path: str):
    wb = Workbook()

    # ── Sheet 1: Raw Data ──────────────────────────────────────────────────────
    ws = wb.active
    ws.title = "Raw Data"
    ws.freeze_panes = "A2"
    for ci, col in enumerate(df.columns, 1):
        ws.cell(row=1, column=ci).value = col
    style_header(ws, 1, len(df.columns))
    for ri, row in df.iterrows():
        for ci, val in enumerate(row, 1):
            cell = ws.cell(row=ri + 2, column=ci)
            cell.value     = val
            cell.alignment = _align(wrap=True)
            cell.border    = _bdr()
            if ri % 2 == 1:
                cell.fill = _fill(OPXL_HEX["light"])
    for col in ws.columns:
        w = max((len(str(c.value or "")) for c in col), default=10)
        ws.column_dimensions[col[0].column_letter].width = min(w + 4, 45)

    # ── Sheet 2: Summary Dashboard ─────────────────────────────────────────────
    ws2 = wb.create_sheet("Summary Dashboard")
    ws2.merge_cells("A1:H2")
    c = ws2["A1"]
    c.value     = "PROJECT LESSONS LEARNED  —  Insights Dashboard"
    c.font      = _font(True, "FFFFFFFF", 16)
    c.fill      = _fill(OPXL_HEX["primary"])
    c.alignment = _align("center", "center")

    tc = col_map.get("project_type")
    sc = col_map.get("project_stage")
    cc = col_map.get("new_category") or col_map.get("original_category")
    kpis = [
        ("Total Records",  len(df),                              OPXL_HEX["primary"]),
        ("Project Types",  df[tc].nunique() if tc else "—",     OPXL_HEX["secondary"]),
        ("Stages",         df[sc].nunique() if sc else "—",     OPXL_HEX["accent1"]),
        ("Categories",     df[cc].nunique() if cc else "—",     OPXL_HEX["accent3"]),
    ]
    for i, (label, val, col_hex) in enumerate(kpis):
        cs = i * 2 + 1
        ws2.merge_cells(start_row=4, start_column=cs, end_row=4, end_column=cs + 1)
        ws2.merge_cells(start_row=5, start_column=cs, end_row=5, end_column=cs + 1)
        c1 = ws2.cell(row=4, column=cs)
        c1.value = label; c1.font = _font(True, "FF606060", 9); c1.alignment = _align("center")
        c2 = ws2.cell(row=5, column=cs)
        c2.value = val; c2.font = _font(True, col_hex, 20)
        c2.alignment = _align("center"); c2.fill = _fill(OPXL_HEX["light"])

    for key, col_l, row_n in [
        ("type_bar",       "A",  7),  ("stage_bar",      "I",  7),
        ("category_bar",   "A", 29),  ("threat_opp_pie", "I", 29),
        ("keyword_all",    "A", 51),  ("wordcloud",      "I", 51),
        ("stage_cat_stack","A", 73),
    ]:
        embed_image(ws2, charts.get(key), f"{col_l}{row_n}", 13, 8)

    for ltr in "ABCDEFGHIJKLMNOP":
        ws2.column_dimensions[ltr].width = 10

    # ── Sheet 3: Frequency Tables ──────────────────────────────────────────────
    ws3  = wb.create_sheet("Frequency Tables")
    crow = 1
    for field in ["project_type","project_stage","new_category",
                  "original_category","organization","threat_opp_type"]:
        col = col_map.get(field)
        if not col:
            continue
        df_f = freq_table(df, col, field.replace("_", " ").title())
        crow = write_freq_block(
            ws3, df_f, crow, f"Distribution by {field.replace('_', ' ').title()}")
    for ltr in ["A","B","C"]:
        ws3.column_dimensions[ltr].width = 35 if ltr == "A" else 14

    # ── Sheet 4: Keyword Analysis ──────────────────────────────────────────────
    ws4 = wb.create_sheet("Keyword Analysis")
    ws4.merge_cells("A1:E1")
    ws4["A1"].value = "Top Keywords by Text Field"
    ws4["A1"].font  = _font(True, OPXL_HEX["primary"], 13)
    kw_row = 3
    for field, label in [
        ("description",    "Description / Threat-Opportunity"),
        ("response",       "Response / Mitigation"),
        ("lessons_learned","Lessons Learned"),
        ("lesson_title",   "Lesson Titles"),
    ]:
        col = col_map.get(field)
        if not col:
            continue
        kws = top_keywords(df[col].tolist(), sw, 20)
        if not kws:
            continue
        ws4.cell(row=kw_row, column=1).value = f">> {label}"
        ws4.cell(row=kw_row, column=1).font  = _font(True, OPXL_HEX["secondary"], 11)
        kw_row += 1
        for ci, h in enumerate(["Keyword", "Frequency"], 1):
            ws4.cell(row=kw_row, column=ci).value = h
        style_header(ws4, kw_row, 2)
        kw_row += 1
        for word, freq in kws:
            ws4.cell(row=kw_row, column=1).value = word
            ws4.cell(row=kw_row, column=2).value = freq
            for ci in [1, 2]:
                ws4.cell(row=kw_row, column=ci).border = _bdr()
                if kw_row % 2 == 0:
                    ws4.cell(row=kw_row, column=ci).fill = _fill(OPXL_HEX["light"])
            kw_row += 1
        kw_row += 2
    ws4.column_dimensions["A"].width = 28
    ws4.column_dimensions["B"].width = 14

    # ── Sheet 5: Threat vs. Opportunity keywords ───────────────────────────────
    ws5 = wb.create_sheet("Threat vs Opportunity")
    ws5.merge_cells("A1:F1")
    ws5["A1"].value = "Threat vs. Opportunity — Keyword Comparison"
    ws5["A1"].font  = _font(True, OPXL_HEX["primary"], 13)
    t_texts, o_texts = split_threats_opps(df, col_map)
    t_kws = top_keywords(t_texts, sw, 15) if t_texts else []
    o_kws = top_keywords(o_texts, sw, 15) if o_texts else []
    row = 3
    ws5.cell(row=row, column=1).value = "THREAT Keywords"
    ws5.cell(row=row, column=1).font  = _font(True, OPXL_HEX["threat"], 11)
    ws5.cell(row=row, column=4).value = "OPPORTUNITY Keywords"
    ws5.cell(row=row, column=4).font  = _font(True, OPXL_HEX["opportunity"], 11)
    row += 1
    for ci, h in [(1,"Keyword"),(2,"Freq"),(4,"Keyword"),(5,"Freq")]:
        ws5.cell(row=row, column=ci).value = h
    style_header(ws5, row, 5, bg=OPXL_HEX["dark"])
    row += 1
    for i in range(max(len(t_kws), len(o_kws))):
        if i < len(t_kws):
            ws5.cell(row=row+i, column=1).value = t_kws[i][0]
            ws5.cell(row=row+i, column=2).value = t_kws[i][1]
        if i < len(o_kws):
            ws5.cell(row=row+i, column=4).value = o_kws[i][0]
            ws5.cell(row=row+i, column=5).value = o_kws[i][1]
        for ci in [1,2,4,5]:
            ws5.cell(row=row+i, column=ci).border = _bdr()
            if i % 2 == 0:
                ws5.cell(row=row+i, column=ci).fill = _fill(OPXL_HEX["light"])
    for ltr in ["A","D"]: ws5.column_dimensions[ltr].width = 28
    for ltr in ["B","E"]: ws5.column_dimensions[ltr].width = 12

    # ── Sheet 6: Stage × Category Heatmap ─────────────────────────────────────
    pivot = stage_cat_pivot(df, col_map)
    if pivot is not None and not pivot.empty:
        ws6 = wb.create_sheet("Stage x Category Matrix")
        ws6["A1"].value = "Project Stage × Category — Heatmap"
        ws6["A1"].font  = _font(True, OPXL_HEX["primary"], 13)
        pr = pivot.reset_index()
        for ci, col in enumerate(pr.columns, 1):
            ws6.cell(row=3, column=ci).value = str(col)
        style_header(ws6, 3, len(pr.columns))
        num_cols = pr.select_dtypes(include="number")
        max_val  = num_cols.max().max() if not num_cols.empty else 1
        for ri, row in pr.iterrows():
            for ci, val in enumerate(row, 1):
                cell = ws6.cell(row=ri + 4, column=ci)
                cell.value     = val
                cell.border    = _bdr()
                cell.alignment = _align("center")
                if ci > 1 and isinstance(val, (int, float)) and val > 0 and max_val:
                    ix = int(val / max_val * 200)
                    cell.fill = _fill(f"FF{format(255-ix,'02X')}F0F0")
        for col in ws6.columns:
            w = max((len(str(c.value or "")) for c in col), default=10)
            ws6.column_dimensions[col[0].column_letter].width = min(w + 4, 30)

    # ── Sheet 7: N-gram Analysis (NEW) ─────────────────────────────────────────
    ngrams = analytics.get("ngrams", {})
    if ngrams:
        ws7 = wb.create_sheet("N-gram Analysis")
        ws7.merge_cells("A1:D1")
        ws7["A1"].value = "N-gram Analysis — Top Bigrams & Trigrams per Text Field"
        ws7["A1"].font  = _font(True, OPXL_HEX["primary"], 13)
        ng_row = 3
        for field, field_ngrams in ngrams.items():
            for ng_type, ng_data in field_ngrams.items():
                if not ng_data:
                    continue
                label = f"{field.replace('_', ' ').title()} — {ng_type.title()}"
                ws7.cell(row=ng_row, column=1).value = f">> {label}"
                ws7.cell(row=ng_row, column=1).font  = _font(True, OPXL_HEX["secondary"], 11)
                ng_row += 1
                for ci, h in enumerate(["Phrase", "Frequency"], 1):
                    ws7.cell(row=ng_row, column=ci).value = h
                style_header(ws7, ng_row, 2)
                ng_row += 1
                for phrase, cnt in ng_data:
                    ws7.cell(row=ng_row, column=1).value = phrase
                    ws7.cell(row=ng_row, column=2).value = cnt
                    for ci in [1, 2]:
                        ws7.cell(row=ng_row, column=ci).border = _bdr()
                        if ng_row % 2 == 0:
                            ws7.cell(row=ng_row, column=ci).fill = _fill(OPXL_HEX["light"])
                    ng_row += 1
                ng_row += 2
        ws7.column_dimensions["A"].width = 45
        ws7.column_dimensions["B"].width = 14

    # ── Sheet 8: Sentiment Analysis (NEW) ─────────────────────────────────────
    sent_df = analytics.get("sentiment")
    lc = col_map.get("lessons_learned")
    if sent_df is not None and lc:
        ws8 = wb.create_sheet("Sentiment Analysis")
        ws8.merge_cells("A1:F1")
        ws8["A1"].value = "Sentiment Analysis — VADER Scores per Record"
        ws8["A1"].font  = _font(True, OPXL_HEX["primary"], 13)

        headers = ["Lesson Text (truncated)", "Sentiment", "Compound", "Positive", "Neutral", "Negative"]
        for ci, h in enumerate(headers, 1):
            ws8.cell(row=3, column=ci).value = h
        style_header(ws8, 3, len(headers))

        for ri in range(len(sent_df)):
            row_data = [
                str(df[lc].iloc[ri])[:120],
                sent_df["label"].iloc[ri],
                round(sent_df["compound"].iloc[ri], 3),
                round(sent_df["pos"].iloc[ri], 3),
                round(sent_df["neu"].iloc[ri], 3),
                round(sent_df["neg"].iloc[ri], 3),
            ]
            bg = OPXL_HEX["light"] if ri % 2 == 0 else "FFFFFFFF"
            for ci, val in enumerate(row_data, 1):
                cell = ws8.cell(row=ri + 4, column=ci)
                cell.value     = val
                cell.fill      = _fill(bg)
                cell.font      = _font(size=9)
                cell.alignment = _align(wrap=True)
                cell.border    = _bdr()
            # Colour-code sentiment label
            lbl_cell = ws8.cell(row=ri + 4, column=2)
            lbl = sent_df["label"].iloc[ri]
            if lbl == "Positive":
                lbl_cell.font = _font(True, OPXL_HEX["opportunity"], 9)
            elif lbl == "Negative":
                lbl_cell.font = _font(True, OPXL_HEX["threat"], 9)

        ws8.column_dimensions["A"].width = 60
        for ltr in ["B","C","D","E","F"]:
            ws8.column_dimensions[ltr].width = 14

        # Embed sentiment charts
        embed_image(ws8, charts.get("sentiment_pie"),  f"H3",  13, 8)
        embed_image(ws8, charts.get("sentiment_hist"), f"H25", 13, 7)

    # ── Sheet 9: Topic Modeling (NEW) ─────────────────────────────────────────
    topics = analytics.get("topics")
    if topics:
        ws9 = wb.create_sheet("Topic Modeling")
        ws9.merge_cells("A1:J1")
        ws9["A1"].value = "LDA Topic Modeling — Top Terms per Topic"
        ws9["A1"].font  = _font(True, OPXL_HEX["primary"], 13)
        for ci, topic_words in enumerate(topics, 1):
            ws9.cell(row=3, column=ci).value = f"Topic {ci}"
            style_header(ws9, 3, len(topics))
            for ri, word in enumerate(topic_words, 1):
                cell = ws9.cell(row=ri + 3, column=ci)
                cell.value     = word
                cell.font      = _font(size=10)
                cell.border    = _bdr()
                cell.alignment = _align("center")
                if ri % 2 == 0:
                    cell.fill = _fill(OPXL_HEX["light"])
        for ci in range(1, len(topics) + 1):
            ws9.column_dimensions[ws9.cell(row=3, column=ci).column_letter].width = 20
        embed_image(ws9, charts.get("topic_model"), "A17", 14 * len(topics) / 5, 8)

    # ── Sheet 10: Advanced Analytics (NEW) ────────────────────────────────────
    ws10 = wb.create_sheet("Advanced Analytics")
    ws10.merge_cells("A1:H1")
    ws10["A1"].value = "Advanced Analytics — TF-IDF, Readability, NER, Clusters"
    ws10["A1"].font  = _fill = _font(True, OPXL_HEX["primary"], 13)
    adv_row = 3

    # TF-IDF per category
    tfidf_cats = analytics.get("tfidf_categories", {})
    if tfidf_cats:
        ws10.cell(row=adv_row, column=1).value = "TF-IDF Distinctive Keywords per Category"
        ws10.cell(row=adv_row, column=1).font  = _font(True, OPXL_HEX["primary"], 12)
        adv_row += 1
        for cat, terms in tfidf_cats.items():
            ws10.cell(row=adv_row, column=1).value = str(cat)
            ws10.cell(row=adv_row, column=1).font  = _font(True, OPXL_HEX["secondary"], 10)
            ws10.cell(row=adv_row, column=2).value = "Term"
            ws10.cell(row=adv_row, column=3).value = "TF-IDF Score"
            style_header(ws10, adv_row, 3, bg=OPXL_HEX["dark"])
            adv_row += 1
            for ri, (word, score) in enumerate(terms):
                bg = OPXL_HEX["light"] if ri % 2 == 0 else "FFFFFFFF"
                for ci, val in enumerate([ri + 1, word, score], 1):
                    cell = ws10.cell(row=adv_row, column=ci)
                    cell.value = val; cell.border = _bdr()
                    cell.fill  = _fill(bg); cell.font = _font(size=9)
                adv_row += 1
            adv_row += 1
        adv_row += 2

    # Readability summary
    read_df = analytics.get("readability")
    if read_df is not None:
        ws10.cell(row=adv_row, column=1).value = "Readability Metrics (Lessons Learned)"
        ws10.cell(row=adv_row, column=1).font  = _font(True, OPXL_HEX["primary"], 12)
        adv_row += 1
        rd_clean = read_df.dropna(subset=["flesch_kincaid"])
        if not rd_clean.empty:
            headers = ["Metric", "Min", "Max", "Mean", "Median"]
            for ci, h in enumerate(headers, 1):
                ws10.cell(row=adv_row, column=ci).value = h
            style_header(ws10, adv_row, len(headers))
            adv_row += 1
            for metric in ["flesch_kincaid", "avg_word_length", "word_count"]:
                col_data = rd_clean[metric].dropna()
                if col_data.empty:
                    continue
                for ci, val in enumerate([
                    metric.replace("_", " ").title(),
                    round(col_data.min(), 2), round(col_data.max(), 2),
                    round(col_data.mean(), 2), round(col_data.median(), 2),
                ], 1):
                    cell = ws10.cell(row=adv_row, column=ci)
                    cell.value = val; cell.border = _bdr()
                    cell.font  = _font(size=9)
                adv_row += 1
        adv_row += 2

    # NER entities
    ner = analytics.get("ner")
    if ner:
        ws10.cell(row=adv_row, column=1).value = "Named Entity Recognition — Top Entities"
        ws10.cell(row=adv_row, column=1).font  = _font(True, OPXL_HEX["primary"], 12)
        adv_row += 1
        for ci, h in enumerate(["Entity", "Type", "Frequency"], 1):
            ws10.cell(row=adv_row, column=ci).value = h
        style_header(ws10, adv_row, 3)
        adv_row += 1
        for i, ((ent, lbl), cnt) in enumerate(ner.most_common(30)):
            bg = OPXL_HEX["light"] if i % 2 == 0 else "FFFFFFFF"
            for ci, val in enumerate([ent, lbl, cnt], 1):
                cell = ws10.cell(row=adv_row, column=ci)
                cell.value = val; cell.border = _bdr()
                cell.fill  = _fill(bg); cell.font = _font(size=9)
            adv_row += 1
        adv_row += 2

    # Embedded charts on this sheet
    embed_image(ws10, charts.get("readability"),    "I3",  13, 7)
    embed_image(ws10, charts.get("ner_chart"),      "I25", 13, 7)
    embed_image(ws10, charts.get("cluster_scatter"),"I47", 13, 8)
    embed_image(ws10, charts.get("cooccurrence"),   "I69", 13, 8)
    embed_image(ws10, charts.get("length_hist"),    "I91", 13, 6)

    for ltr in ["A","B","C","D"]:
        ws10.column_dimensions[ltr].width = 30 if ltr == "A" else 20

    wb.save(out_path)
    print(f"  Excel saved  ->  {out_path}")


# =============================================================================
# POWERPOINT BUILDER
# =============================================================================

W = Inches(13.33)
H = Inches(7.5)


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color


def _rect(slide, l, t, w, h, color: RGBColor):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s


def _textbox(slide, text, l, t, w, h, size=14, bold=False,
             color=None, align=PP_ALIGN.LEFT):
    color = color or RGBColor(0x1A, 0x2B, 0x3C)
    tb    = slide.shapes.add_textbox(l, t, w, h)
    tf    = tb.text_frame; tf.word_wrap = True
    p     = tf.paragraphs[0]; p.alignment = align
    run   = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color


def slide_header(prs, slide, title, subtitle=""):
    _bg(slide, RGBColor(0xFA, 0xFA, 0xFB))
    _rect(slide, 0, 0, W, Inches(0.85), PPTX_RGB["primary"])
    _textbox(slide, title, Inches(0.3), Inches(0.08), Inches(12.5), Inches(0.65),
             size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        _textbox(slide, subtitle, Inches(0.3), Inches(0.72), Inches(12), Inches(0.3),
                 size=10, color=RGBColor(0xBB, 0xCC, 0xDD))


def add_image_slide(prs, path, title, subtitle=""):
    s = _slide(prs)
    slide_header(prs, s, title, subtitle)
    if path and Path(path).exists():
        s.shapes.add_picture(path, Inches(0.5), Inches(1.05),
                             width=Inches(12.3), height=Inches(6.1))


def add_two_image_slide(prs, path_l, path_r, title, subtitle=""):
    """Two charts side-by-side on one slide."""
    s = _slide(prs)
    slide_header(prs, s, title, subtitle)
    for path, left in [(path_l, Inches(0.3)), (path_r, Inches(6.8))]:
        if path and Path(path).exists():
            s.shapes.add_picture(path, left, Inches(1.05),
                                 width=Inches(6.2), height=Inches(6.0))


def add_table_slide(prs, df_t, title, col_widths=None):
    s = _slide(prs)
    slide_header(prs, s, title)
    rows, cols = len(df_t) + 1, len(df_t.columns)
    tbl = s.shapes.add_table(
        rows, cols, Inches(0.4), Inches(1.0),
        Inches(12.5), Inches(min(5.8, rows * 0.38))).table
    col_widths = col_widths or [Inches(12.5 / cols)] * cols
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw
    for ci, col in enumerate(df_t.columns):
        cell = tbl.cell(0, ci); cell.text = str(col)
        cell.fill.solid(); cell.fill.fore_color.rgb = PPTX_RGB["primary"]
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.bold = True; p.runs[0].font.size = Pt(10)
        p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for ri, row in df_t.iterrows():
        bg = RGBColor(0xF0, 0xF4, 0xF8) if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci); cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(9)
            p.runs[0].font.color.rgb = RGBColor(0x1A, 0x2B, 0x3C)


def build_pptx(df: pd.DataFrame, col_map: dict, charts: dict,
               sw: set, analytics: dict, out_path: str):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Title slide
    s = _slide(prs)
    _bg(s, PPTX_RGB["dark"])
    _rect(s, 0, 0,           W, Inches(0.1), PPTX_RGB["accent1"])
    _rect(s, 0, Inches(7.4), W, Inches(0.1), PPTX_RGB["accent1"])
    _textbox(s, "PROJECT LESSONS LEARNED",
             Inches(1), Inches(2.0), Inches(11), Inches(1.1),
             size=40, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    _textbox(s, "Text Analytics & Insight Report",
             Inches(1), Inches(3.2), Inches(11), Inches(0.7),
             size=20, color=PPTX_RGB["accent2"], align=PP_ALIGN.CENTER)
    _textbox(s, f"Total Records Analysed: {len(df):,}",
             Inches(1), Inches(4.0), Inches(11), Inches(0.5),
             size=13, color=RGBColor(0x88,0xAA,0xCC), align=PP_ALIGN.CENTER)

    # Agenda
    s2 = _slide(prs)
    slide_header(prs, s2, "Agenda")
    items = [
        "01  Distribution by Project Type",
        "02  Distribution by Stage & Category",
        "03  Threat vs. Opportunity Split",
        "04  Top Keywords (All Text / Lessons)",
        "05  Word Cloud — Key Themes",
        "06  Stage × Category Matrix",
        "07  N-gram Analysis (Bigrams & Trigrams)",
        "08  Sentiment Analysis",
        "09  LDA Topic Modeling",
        "10  Document Clusters (PCA)",
        "11  Sankey Flow Diagram",
        "12  Word Co-occurrence Network",
        "13  Key Takeaways",
    ]
    for i, item in enumerate(items):
        y = Inches(1.1 + i * 0.47)
        _rect(s2, Inches(0.4), y + Inches(0.06), Inches(0.05), Inches(0.35), PPTX_RGB["accent1"])
        _textbox(s2, item, Inches(0.6), y, Inches(11), Inches(0.45), size=12)

    # ── Original chart slides ─────────────────────────────────────────────────
    for key, title, subtitle in [
        ("type_bar",       "Distribution by Project Type",     "Lessons by project type"),
        ("stage_bar",      "Distribution by Project Stage",    "Which stages generate the most lessons"),
        ("category_bar",   "Distribution by Category",         "Risk/opportunity category breakdown"),
        ("threat_opp_pie", "Threat vs. Opportunity Split",     "Proportion of threats and opportunities"),
        ("keyword_all",    "Top Keywords — All Text",          "Most frequent terms across all text fields"),
        ("keyword_lessons","Top Keywords — Lessons Learned",   "Key terms from lessons learned text"),
        ("threat_kw",      "Top Keywords — Threats",           "Dominant threat-related themes"),
        ("opp_kw",         "Top Keywords — Opportunities",     "Dominant opportunity-related themes"),
        ("wordcloud",      "Word Cloud — Combined Text",       "Visual map of prominent themes"),
        ("stage_cat_stack","Stage × Category Breakdown",       "How categories distribute across stages"),
    ]:
        path = charts.get(key)
        if path and Path(path).exists():
            add_image_slide(prs, path, title, subtitle)

    # Frequency table slides
    for field in ["project_type", "project_stage", "new_category"]:
        col = col_map.get(field)
        if not col:
            continue
        df_f = freq_table(df, col, field.replace("_", " ").title())
        add_table_slide(prs, df_f.head(10),
                        f"Top 10 — {field.replace('_', ' ').title()}",
                        col_widths=[Inches(7), Inches(2.5), Inches(3)])

    pivot = stage_cat_pivot(df, col_map)
    if pivot is not None and not pivot.empty:
        pr  = pivot.reset_index()
        nc  = len(pr.columns)
        cws = [Inches(2.5)] + [Inches(10 / max(nc - 1, 1))] * (nc - 1)
        add_table_slide(prs, pr.head(10), "Stage × Category Matrix", col_widths=cws)

    # ── New v3 chart slides ───────────────────────────────────────────────────

    # N-gram: pair bigrams + trigrams side by side for each field
    ngrams = analytics.get("ngrams", {})
    for field in ["description", "lessons_learned"]:
        field_ngrams = ngrams.get(field, {})
        bg_key = f"{field}_bigrams"
        tg_key = f"{field}_trigrams"
        if charts.get(bg_key) or charts.get(tg_key):
            label = field.replace("_", " ").title()
            add_two_image_slide(prs,
                                charts.get(bg_key), charts.get(tg_key),
                                f"N-gram Analysis — {label}",
                                "Bigrams (left) and Trigrams (right)")

    # Sentiment charts (side by side)
    if charts.get("sentiment_pie") or charts.get("sentiment_hist"):
        add_two_image_slide(prs,
                            charts.get("sentiment_pie"), charts.get("sentiment_hist"),
                            "Sentiment Analysis (VADER)",
                            "Distribution (left) and compound score histogram (right)")

    # Text length histogram
    if charts.get("length_hist"):
        add_image_slide(prs, charts["length_hist"],
                        "Lesson Length Distribution",
                        "Word count per record — identifies outliers and data quality issues")

    # Readability chart
    if charts.get("readability"):
        add_image_slide(prs, charts["readability"],
                        "Readability by Stage / Category",
                        "Flesch-Kincaid grade level — lower = easier to read")

    # LDA topic modeling
    if charts.get("topic_model"):
        add_image_slide(prs, charts["topic_model"],
                        "LDA Topic Modeling",
                        "Latent topics discovered across all lessons learned text")

    # Document cluster scatter
    if charts.get("cluster_scatter"):
        add_image_slide(prs, charts["cluster_scatter"],
                        "Document Clusters (TF-IDF + K-Means + PCA)",
                        "Each dot = one record; colour = cluster. Shows natural groupings.")

    # Sankey
    if charts.get("sankey"):
        add_image_slide(prs, charts["sankey"],
                        "Flow Diagram: Type → Stage → Category",
                        "Sankey diagram showing how records flow across dimensions")

    # Co-occurrence network
    if charts.get("cooccurrence"):
        add_image_slide(prs, charts["cooccurrence"],
                        "Word Co-occurrence Network",
                        "Nodes = keywords; edges = how often they appear in the same record")

    # NER chart
    if charts.get("ner_chart"):
        add_image_slide(prs, charts["ner_chart"],
                        "Top Named Entities (NER)",
                        "People, organizations, and locations most frequently mentioned")

    # Similarity heatmap
    if charts.get("similarity"):
        add_image_slide(prs, charts["similarity"],
                        "Semantic Similarity Matrix",
                        "Cosine similarity between lesson embeddings — darker = more similar")

    # Key takeaways slide
    s_kt = _slide(prs)
    _bg(s_kt, PPTX_RGB["dark"])
    _rect(s_kt, 0, 0, W, Inches(1.0), PPTX_RGB["accent1"])
    _textbox(s_kt, "Key Takeaways",
             Inches(0.3), Inches(0.1), Inches(12), Inches(0.8),
             size=28, bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    tc_col = col_map.get("project_type")
    sc_col = col_map.get("project_stage")
    cc_col = col_map.get("new_category") or col_map.get("original_category")
    dc_col = col_map.get("description")
    lc_col = col_map.get("lessons_learned")
    t_texts2, o_texts2 = split_threats_opps(df, col_map)

    takeaways = []
    if tc_col:
        top = df[tc_col].value_counts()
        takeaways.append(f"Most lessons come from '{top.idxmax()}' projects ({top.iloc[0]} records)")
    if sc_col:
        top = df[sc_col].value_counts()
        takeaways.append(f"'{top.idxmax()}' is the most represented project stage")
    if cc_col:
        top = df[cc_col].value_counts()
        takeaways.append(f"'{top.idxmax()}' is the leading lesson category")
    tn, on = len(t_texts2), len(o_texts2)
    if tn + on > 0:
        takeaways.append(f"Dataset: {tn} threats, {on} opportunities "
                         f"({tn/(tn+on)*100:.0f}% threats)")
    sent_df = analytics.get("sentiment")
    if sent_df is not None:
        lbl_counts = sent_df["label"].value_counts()
        dominant   = lbl_counts.idxmax()
        pct        = lbl_counts.max() / len(sent_df) * 100
        takeaways.append(f"Sentiment: {dominant.lower()} tone dominates ({pct:.0f}% of records)")
    topics = analytics.get("topics")
    if topics:
        takeaways.append(f"LDA identified {len(topics)} latent topic clusters in the lessons text")
    if dc_col:
        kws = top_keywords(df[dc_col].tolist(), sw, 5)
        if kws:
            takeaways.append("Top description keywords: " + ", ".join(k for k, _ in kws))
    if lc_col:
        kws = top_keywords(df[lc_col].tolist(), sw, 5)
        if kws:
            takeaways.append("Top lessons learned keywords: " + ", ".join(k for k, _ in kws))
    takeaways.append("Full analysis available in the companion Excel workbook")

    for i, ta in enumerate(takeaways[:8]):
        y = Inches(1.15 + i * 0.75)
        _rect(s_kt, Inches(0.4), y + Inches(0.1), Inches(0.35), Inches(0.38), PPTX_RGB["accent2"])
        _textbox(s_kt, str(i + 1), Inches(0.43), y + Inches(0.08),
                 Inches(0.3), Inches(0.38),
                 size=12, bold=True, color=PPTX_RGB["dark"], align=PP_ALIGN.CENTER)
        _textbox(s_kt, ta, Inches(0.9), y, Inches(11.8), Inches(0.65),
                 size=13, color=RGBColor(0xEC, 0xE2, 0xD0))

    prs.save(out_path)
    print(f"  PowerPoint saved  ->  {out_path}")


# =============================================================================
# MAIN
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Analyse project lessons-learned Excel data → insights workbook + PowerPoint")
    parser.add_argument("input_file", help="Input .xlsx file")
    parser.add_argument("--sheet",   default=None, help="Sheet name (default: first sheet)")
    parser.add_argument("--out-dir", default=".",  help="Output directory (default: .)")
    args = parser.parse_args()

    input_path = Path(args.input_file)
    if not input_path.exists():
        sys.exit(f"File not found: {input_path}")

    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    stem = input_path.stem

    print(f"\nLoading:  {input_path}")
    df      = load_data(str(input_path), args.sheet)
    col_map = map_columns(df)
    sw      = build_stopwords()

    print(f"  Rows: {len(df)}  |  Columns: {len(df.columns)}")
    print(f"  Recognised fields: {list(col_map.keys())}\n")

    # Optional packages status
    print("Optional packages:")
    print(f"  nltk={NLTK_OK}  wordcloud={WC_OK}  sklearn={SKLEARN_OK}  "
          f"textstat={TEXTSTAT_OK}  networkx={NX_OK}")
    print(f"  spacy={SPACY_OK}  plotly={PLOTLY_OK}  sentence-transformers={SBERT_OK}\n")

    print("Running NLP analytics …")
    analytics = run_analytics(df, col_map, sw)
    print()

    print("Generating charts …")
    charts = generate_charts(df, col_map, sw, analytics)
    print()

    xl_out   = out_dir / f"{stem}_insights.xlsx"
    pptx_out = out_dir / f"{stem}_insights.pptx"

    print("Building Excel workbook …")
    build_excel(df, col_map, charts, sw, analytics, str(xl_out))

    print("Building PowerPoint presentation …")
    build_pptx(df, col_map, charts, sw, analytics, str(pptx_out))

    print(f"\nDone!  Outputs in: {out_dir.resolve()}")
    print(f"  {xl_out.name}")
    print(f"  {pptx_out.name}\n")


if __name__ == "__main__":
    main()
